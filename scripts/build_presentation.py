#!/usr/bin/env python3
"""Sunum üretim aracı — markdown kaynaktan PPTX oluşturur.

Kullanım:
    python scripts/build_presentation.py
    python scripts/build_presentation.py --girdi presentations/final_sunumu.md --cikti presentations/final_sunumu.pptx

Girdi biçimi (presentations/*.md):
    - Slaytlar `---` satırıyla ayrılır.
    - Her slaydın ilk `# ` satırı slayt başlığıdır.
    - `- ` ile başlayan satırlar madde imidir; iki boşluk girintili
      `  - ` satırlar alt madde olur.
    - `> not:` ile başlayan satırlar konuşmacı notuna yazılır.
    - Diğer düz satırlar normal metin paragrafı olarak eklenir.

Not: PPTX çıktısının PDF sürümü (şartname iki formatı da ister) PowerPoint
üzerinden "Farklı Kaydet → PDF" ile elle alınmalıdır.
"""

from __future__ import annotations

import argparse
import sys
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.util import Emu, Inches, Pt
except ImportError:  # pragma: no cover - kurulum yönlendirmesi
    print(
        "HATA: python-pptx kurulu değil.\n"
        "Kurulum: pip install -r requirements-optional.txt",
        file=sys.stderr,
    )
    sys.exit(1)

PROJE_KOKU = Path(__file__).resolve().parent.parent

# Kurumsal renkler (TEKNOFEST tonlarına yakın koyu lacivert + kırmızı vurgu)
RENK_KOYU = RGBColor(0x1B, 0x2A, 0x4A)
RENK_VURGU = RGBColor(0xC8, 0x10, 0x2E)
RENK_METIN = RGBColor(0x21, 0x21, 0x21)
RENK_BEYAZ = RGBColor(0xFF, 0xFF, 0xFF)

YAZI_TIPI = "Calibri"


def _md_ayristir(md_metin: str) -> list[dict]:
    """Markdown metnini slayt sözlükleri listesine çevirir.

    Her slayt: {"baslik": str, "maddeler": [(seviye, metin)], "notlar": [str]}
    """
    slaytlar = []
    for blok in md_metin.split("\n---\n"):
        blok = blok.strip()
        if not blok:
            continue
        slayt = {"baslik": "", "maddeler": [], "notlar": []}
        for satir in blok.splitlines():
            ham = satir.rstrip()
            if not ham.strip():
                continue
            if ham.startswith("# ") and not slayt["baslik"]:
                slayt["baslik"] = ham[2:].strip()
            elif ham.lstrip().startswith("> not:"):
                slayt["notlar"].append(ham.lstrip()[6:].strip())
            elif ham.startswith("  - "):
                slayt["maddeler"].append((1, ham[4:].strip()))
            elif ham.startswith("- "):
                slayt["maddeler"].append((0, ham[2:].strip()))
            else:
                # Düz paragraf — seviye 0 madde olarak eklenir (imsiz)
                slayt["maddeler"].append((0, ham.strip()))
        if slayt["baslik"] or slayt["maddeler"]:
            slaytlar.append(slayt)
    return slaytlar


def _vurgu_bandi(slide, prs) -> None:
    """Slaydın sol kenarına ince dikey vurgu bandı ekler."""
    from pptx.enum.shapes import MSO_SHAPE

    bant = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Emu(0), Emu(0), Inches(0.18), prs.slide_height
    )
    bant.fill.solid()
    bant.fill.fore_color.rgb = RENK_VURGU
    bant.line.fill.background()


def _kapak_slaydi(prs, slayt: dict) -> None:
    """İlk slaydı koyu zeminli kapak olarak biçimlendirir."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # boş yerleşim

    zemin = slide.background
    zemin.fill.solid()
    zemin.fill.fore_color.rgb = RENK_KOYU

    _vurgu_bandi(slide, prs)

    baslik = slide.shapes.add_textbox(
        Inches(0.9), Inches(1.6), Inches(11.5), Inches(2.4)
    )
    tf = baslik.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = slayt["baslik"]
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = RENK_BEYAZ
    p.font.name = YAZI_TIPI

    alt = slide.shapes.add_textbox(Inches(0.9), Inches(4.2), Inches(11.5), Inches(2.6))
    tf = alt.text_frame
    tf.word_wrap = True
    for i, (_, metin) in enumerate(slayt["maddeler"]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = metin
        p.font.size = Pt(20)
        p.font.color.rgb = RENK_BEYAZ
        p.font.name = YAZI_TIPI
        p.space_after = Pt(6)

    _notlari_yaz(slide, slayt)


def _icerik_slaydi(prs, slayt: dict) -> None:
    """Standart içerik slaydı: başlık + madde listesi."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # boş yerleşim
    _vurgu_bandi(slide, prs)

    baslik = slide.shapes.add_textbox(
        Inches(0.7), Inches(0.35), Inches(12.0), Inches(1.0)
    )
    tf = baslik.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = slayt["baslik"]
    p.font.size = Pt(30)
    p.font.bold = True
    p.font.color.rgb = RENK_KOYU
    p.font.name = YAZI_TIPI

    govde = slide.shapes.add_textbox(
        Inches(0.8), Inches(1.45), Inches(12.0), Inches(5.7)
    )
    tf = govde.text_frame
    tf.word_wrap = True
    # Madde sayısına göre yazı boyutunu kademelendir (taşmayı önler)
    madde_sayisi = len(slayt["maddeler"])
    boyut = Pt(18) if madde_sayisi <= 8 else Pt(16) if madde_sayisi <= 11 else Pt(14)
    for i, (seviye, metin) in enumerate(slayt["maddeler"]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = ("• " if seviye == 0 else "– ") + metin
        p.level = seviye
        p.font.size = boyut if seviye == 0 else Pt(max(boyut.pt - 2, 12))
        p.font.color.rgb = RENK_METIN
        p.font.name = YAZI_TIPI
        p.space_after = Pt(5)

    _notlari_yaz(slide, slayt)


def _notlari_yaz(slide, slayt: dict) -> None:
    """Konuşmacı notlarını slaydın not bölümüne ekler."""
    if slayt["notlar"]:
        slide.notes_slide.notes_text_frame.text = "\n".join(slayt["notlar"])


def sunum_uret(girdi: Path, cikti: Path) -> int:
    """Markdown kaynaktan PPTX üretir; slayt sayısını döndürür."""
    slaytlar = _md_ayristir(girdi.read_text(encoding="utf-8"))
    if not slaytlar:
        print(f"HATA: {girdi} içinde slayt bulunamadı.", file=sys.stderr)
        sys.exit(1)

    prs = Presentation()
    prs.slide_width = Inches(13.333)  # 16:9
    prs.slide_height = Inches(7.5)

    _kapak_slaydi(prs, slaytlar[0])
    for slayt in slaytlar[1:]:
        _icerik_slaydi(prs, slayt)

    cikti.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(cikti))
    return len(slaytlar)


def main() -> None:
    parser = argparse.ArgumentParser(
        description="Markdown kaynaktan sunum (PPTX) üretir."
    )
    parser.add_argument(
        "--girdi",
        type=Path,
        default=PROJE_KOKU / "presentations" / "on_degerlendirme_sunumu.md",
        help="Slayt kaynağı markdown dosyası",
    )
    parser.add_argument(
        "--cikti",
        type=Path,
        default=PROJE_KOKU / "presentations" / "on_degerlendirme_sunumu.pptx",
        help="Üretilecek PPTX dosyası",
    )
    args = parser.parse_args()

    adet = sunum_uret(args.girdi, args.cikti)
    # Windows konsolları UTF-8 olmayabilir; çıktı konsol-güvenli tutulur
    print(f"Sunum uretildi: {args.cikti} ({adet} slayt)")
    print("Hatirlatma: PDF surumunu PowerPoint'ten 'Farkli Kaydet -> PDF' ile alin.")


if __name__ == "__main__":
    main()
