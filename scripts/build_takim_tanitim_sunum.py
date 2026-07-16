#!/usr/bin/env python3
# Copyright 2026 AGENTRA TECH
# SPDX-License-Identifier: Apache-2.0

"""AGENTRA TECH — Takım Tanıtım Sunumu üreticisi (görsel/diyagramlı/kod bloklu).

Bu script, TEKNOFEST 2026 Yapay Zeka Dil Ajanları Yarışması (1. Senaryo) için
kurumsal-profesyonel bir Takım Tanıtım Sunumu (PPTX) üretir. Slaytlar; native
PowerPoint şekilleriyle çizilmiş mimari diyagramlar, koşullu kapı akışı, KPI
kartları, üye kart ızgarası, metrik çubukları, yol haritası ve sözdizimi-renkli
gerçek kod blokları içerir. Metin bilinçli olarak azdır; anlatı görselle taşınır.

Kullanım:
    pip install -r requirements-optional.txt   # python-pptx (yalnızca ilk sefer)
    python scripts/build_takim_tanitim_sunum.py
    # PDF sürümü: PowerPoint'te açıp Dosya -> Farklı Kaydet -> PDF

İçerik kaynakları (dürüstlük): takım/rol bilgileri AUTHORS + takım tanıtım
dosyası; metrikler data/processed/eval_report*.json; kod parçaları gerçek
kaynaktan (src/agents/orchestrator.py, classification_agent.py) sadeleştirilerek.
"""

from __future__ import annotations

import argparse
import re
import sys
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
    from pptx.util import Emu, Inches, Pt
except ImportError:  # pragma: no cover
    print(
        "HATA: python-pptx kurulu değil.\n"
        "Kurulum: pip install -r requirements-optional.txt",
        file=sys.stderr,
    )
    sys.exit(1)

PROJE_KOKU = Path(__file__).resolve().parent.parent

# ----------------------------------------------------------------------------
# Tasarım sistemi — renkler
# ----------------------------------------------------------------------------
def C(hexstr: str) -> RGBColor:
    return RGBColor(int(hexstr[0:2], 16), int(hexstr[2:4], 16), int(hexstr[4:6], 16))


NAVY = C("0E1B33")      # kapak zemini / koyu
NAVY2 = C("16264A")
INK = C("1A2233")       # ana metin
MUTED = C("5B6B85")     # ikincil metin
LINE = C("DCE3EF")      # kart kenarı / ince çizgi
WHITE = C("FFFFFF")
SOFT = C("F3F6FB")      # panel zemini
RED = C("C8102E")       # TEKNOFEST vurgu
BLUE = C("1B6FC2")      # Görev 1 / geliştirme
ORANGE = C("E8590C")    # Görev 2 / adversarial
PURPLE = C("6C3FD1")    # orkestratör
TEAL = C("0CA678")
GREEN = C("2F9E44")     # çıktı / başarı
YELLOW = C("F1B434")    # kapılar
DARKYEL = C("8A5A00")

# Üye renkleri (görev dağılımı lejantıyla tutarlı)
SEYMA = C("1B6FC2")
SINA = C("6C3FD1")
ZEYNEP = C("0CA678")
EMINE = C("E8590C")

# Kod bloğu paleti
CODE_BG = C("0E1626")
CODE_BAR = C("14203A")
CODE_BORDER = C("263B57")
CODE_TEXT = C("D6E2F0")
CODE_KW = C("6FB3FF")
CODE_STR = C("FFB570")
CODE_COM = C("6E8BA6")
CODE_NUM = C("7BE0C3")
CODE_FN = C("E8C07D")

BODY = "Segoe UI"
SEMI = "Segoe UI Semibold"
MONO = "Consolas"

EMU_IN = 914400
W_IN, H_IN = 13.333, 7.5


# ----------------------------------------------------------------------------
# Düşük seviye yardımcılar
# ----------------------------------------------------------------------------
def R(text, color=INK, bold=False, size=14, font=BODY):
    """Metin koşusu (run) tanımı."""
    return (text, color, bold, size, font)


def P(runs, align=PP_ALIGN.LEFT, space_after=2, space_before=0, line=None):
    return {"runs": runs, "align": align, "space_after": space_after,
            "space_before": space_before, "line": line}


def write(tf, paras, anchor=MSO_ANCHOR.TOP, mL=6, mR=6, mT=4, mB=4):
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Pt(mL)
    tf.margin_right = Pt(mR)
    tf.margin_top = Pt(mT)
    tf.margin_bottom = Pt(mB)
    for i, pa in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = pa.get("align", PP_ALIGN.LEFT)
        p.space_after = Pt(pa.get("space_after", 2))
        p.space_before = Pt(pa.get("space_before", 0))
        if pa.get("line"):
            p.line_spacing = pa["line"]
        for (txt, color, bold, size, font) in pa["runs"]:
            r = p.add_run()
            r.text = txt
            r.font.size = Pt(size)
            r.font.bold = bold
            r.font.color.rgb = color
            r.font.name = font
    return tf


def _noshadow(shape):
    try:
        shape.shadow.inherit = False
    except Exception:
        pass


def rrect(slide, x, y, w, h, fill, line=None, line_w=1.0, radius=0.09):
    sp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                Inches(x), Inches(y), Inches(w), Inches(h))
    try:
        sp.adjustments[0] = radius
    except Exception:
        pass
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(line_w)
    _noshadow(sp)
    return sp


def rect(slide, x, y, w, h, fill, line=None, line_w=1.0):
    sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(line_w)
    _noshadow(sp)
    return sp


def oval(slide, x, y, w, h, fill, line=None, line_w=1.0):
    sp = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(line_w)
    _noshadow(sp)
    return sp


def shape_text(sp, paras, anchor=MSO_ANCHOR.MIDDLE, mL=6, mR=6, mT=3, mB=3):
    write(sp.text_frame, paras, anchor=anchor, mL=mL, mR=mR, mT=mT, mB=mB)


def textbox(slide, x, y, w, h, paras, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    write(tb.text_frame, paras, anchor=anchor, mL=0, mR=0, mT=0, mB=0)
    return tb


def arrow(slide, x, y, w, h, color=C("9AA7BD"), direction="right"):
    sh = MSO_SHAPE.RIGHT_ARROW if direction == "right" else MSO_SHAPE.DOWN_ARROW
    sp = slide.shapes.add_shape(sh, Inches(x), Inches(y), Inches(w), Inches(h))
    sp.fill.solid()
    sp.fill.fore_color.rgb = color
    sp.line.fill.background()
    try:
        sp.adjustments[0] = 0.55
        sp.adjustments[1] = 0.55
    except Exception:
        pass
    _noshadow(sp)
    return sp


def node(slide, x, y, w, h, title, sub=None, fill=WHITE, tcolor=INK,
         line=LINE, tsize=12.5, ssize=9, radius=0.12, bold=True):
    sp = rrect(slide, x, y, w, h, fill, line=line, line_w=1.25, radius=radius)
    paras = [P([R(title, tcolor, bold, tsize, SEMI)], align=PP_ALIGN.CENTER,
               space_after=1)]
    if sub:
        paras.append(P([R(sub, tcolor, False, ssize, BODY)],
                       align=PP_ALIGN.CENTER, space_after=0))
    shape_text(sp, paras, anchor=MSO_ANCHOR.MIDDLE)
    return sp


def pill(slide, x, y, w, h, text, fill, tcolor=WHITE, size=9.5, bold=True):
    sp = rrect(slide, x, y, w, h, fill, line=None, radius=0.5)
    shape_text(sp, [P([R(text, tcolor, bold, size, SEMI)], align=PP_ALIGN.CENTER,
                      space_after=0)], anchor=MSO_ANCHOR.MIDDLE)
    return sp


def diamond(slide, x, y, w, h, lines, fill=YELLOW, tcolor=C("3A2B00")):
    sp = slide.shapes.add_shape(MSO_SHAPE.DIAMOND,
                                Inches(x), Inches(y), Inches(w), Inches(h))
    sp.fill.solid()
    sp.fill.fore_color.rgb = fill
    sp.line.color.rgb = DARKYEL
    sp.line.width = Pt(1.0)
    _noshadow(sp)
    paras = [P([R(t, tcolor, b, s, SEMI)], align=PP_ALIGN.CENTER, space_after=0)
             for (t, b, s) in lines]
    shape_text(sp, paras, anchor=MSO_ANCHOR.MIDDLE, mL=2, mR=2)
    return sp


# ----------------------------------------------------------------------------
# Sayfa iskeleti — kapak / başlık / altbilgi
# ----------------------------------------------------------------------------
def new_slide(prs, dark=False):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = NAVY if dark else WHITE
    return slide


def header(slide, kicker, title, accent=RED):
    textbox(slide, 0.62, 0.40, 12.1, 0.32,
            [P([R(kicker.upper(), accent, True, 11, SEMI)])])
    textbox(slide, 0.58, 0.66, 12.2, 0.7,
            [P([R(title, INK, True, 25, SEMI)])])
    rect(slide, 0.64, 1.42, 0.9, 0.055, accent)


def footer(slide, n, total=13):
    rect(slide, 0.62, 7.02, 12.1, 0.014, LINE)
    textbox(slide, 0.62, 7.08, 8.0, 0.3,
            [P([R("AGENTRA TECH  ·  TEKNOFEST 2026 · Yapay Zeka Dil Ajanları Yarışması · 1. Senaryo",
                 MUTED, False, 8.5, BODY)])])
    textbox(slide, 11.3, 7.08, 1.42, 0.3,
            [P([R(f"{n:02d} / {total:02d}", MUTED, False, 8.5, BODY)],
               align=PP_ALIGN.RIGHT)])


# ----------------------------------------------------------------------------
# Kod bloğu (sözdizimi renkli)
# ----------------------------------------------------------------------------
_KW = {"if", "not", "and", "or", "else", "elif", "def", "return", "self",
       "import", "from", "for", "in", "while", "None", "True", "False",
       "class", "try", "except", "is"}
_TOK = re.compile(r"(#.*)|(\"[^\"]*\"|'[^']*')|(\d+\.?\d*)|(\w+)|(\s+)|(.)")


def code_runs(line, size=12.5):
    runs = []
    prev_word = None
    for m in _TOK.finditer(line):
        com, s, num, word, ws, other = m.groups()
        if com is not None:
            runs.append(R(com, CODE_COM, False, size, MONO))
        elif s is not None:
            runs.append(R(s, CODE_STR, False, size, MONO))
        elif num is not None:
            runs.append(R(num, CODE_NUM, False, size, MONO))
        elif word is not None:
            if word in _KW:
                runs.append(R(word, CODE_KW, True, size, MONO))
            elif prev_word == "def":
                runs.append(R(word, CODE_FN, False, size, MONO))
            else:
                runs.append(R(word, CODE_TEXT, False, size, MONO))
            prev_word = word
            continue
        elif ws is not None:
            runs.append(R(ws, CODE_TEXT, False, size, MONO))
        else:
            runs.append(R(other, CODE_TEXT, False, size, MONO))
        prev_word = None
    if not runs:
        runs = [R(" ", CODE_TEXT, False, size, MONO)]
    return runs


def code_block(slide, x, y, w, h, filename, lines, size=12.5):
    rrect(slide, x, y, w, h, CODE_BG, line=CODE_BORDER, line_w=1.0, radius=0.05)
    # başlık şeridi
    rrect(slide, x, y, w, 0.34, CODE_BAR, line=None, radius=0.05)
    for i, dot in enumerate((C("FF5F56"), C("FFBD2E"), C("27C93F"))):
        oval(slide, x + 0.16 + i * 0.2, y + 0.11, 0.11, 0.11, dot)
    textbox(slide, x + 0.78, y + 0.06, w - 1.0, 0.26,
            [P([R(filename, C("8FA6C2"), False, 9.5, MONO)])])
    # kod gövdesi
    tb = slide.shapes.add_textbox(Inches(x + 0.18), Inches(y + 0.42),
                                  Inches(w - 0.36), Inches(h - 0.5))
    tf = tb.text_frame
    tf.word_wrap = False
    tf.margin_left = Pt(2)
    tf.margin_right = Pt(2)
    tf.margin_top = Pt(0)
    tf.margin_bottom = Pt(0)
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(1.5)
        p.line_spacing = 1.0
        for (txt, color, bold, sz, font) in code_runs(ln, size):
            rn = p.add_run()
            rn.text = txt
            rn.font.size = Pt(sz)
            rn.font.bold = bold
            rn.font.color.rgb = color
            rn.font.name = font
    return tb


# ----------------------------------------------------------------------------
# Bileşenler
# ----------------------------------------------------------------------------
def kpi_tile(slide, x, y, w, h, big, unit, label, accent=BLUE):
    rrect(slide, x, y, w, h, WHITE, line=LINE, line_w=1.25, radius=0.1)
    rect(slide, x, y + 0.14, 0.07, h - 0.28, accent)
    runs = [R(big, accent, True, 30, SEMI)]
    if unit:
        runs.append(R(" " + unit, accent, True, 13, SEMI))
    textbox(slide, x + 0.22, y + 0.18, w - 0.3, 0.6, [P(runs)])
    textbox(slide, x + 0.24, y + h - 0.72, w - 0.36, 0.62,
            [P([R(label, INK, False, 10.5, BODY)], line=1.0)])


def card(slide, x, y, w, h, icon, title, body_lines, accent=BLUE):
    rrect(slide, x, y, w, h, WHITE, line=LINE, line_w=1.25, radius=0.09)
    rect(slide, x, y, w, 0.09, accent)
    ttl = [R((icon + "  " if icon else ""), INK, True, 13, SEMI),
           R(title, INK, True, 13.5, SEMI)]
    textbox(slide, x + 0.22, y + 0.24, w - 0.4, 0.5, [P(ttl)])
    paras = [P([R("•  ", accent, True, 10.5, BODY), R(t, INK, False, 10.5, BODY)],
               space_after=3, line=1.02) for t in body_lines]
    textbox(slide, x + 0.22, y + 0.74, w - 0.42, h - 0.9, paras)


def metric_bar(slide, x, y, w, label, value, accent, vtext=None):
    textbox(slide, x, y, w, 0.26, [P([R(label, INK, False, 11, BODY)])])
    track_y = y + 0.28
    rrect(slide, x, track_y, w, 0.2, C("E9EEF6"), line=None, radius=0.5)
    fillw = max(0.12, w * value)
    rrect(slide, x, track_y, fillw, 0.2, accent, line=None, radius=0.5)
    textbox(slide, x + w - 1.2, y, 1.2, 0.26,
            [P([R(vtext or f"{value:.2f}".replace('.', ','), accent, True, 11, SEMI)],
               align=PP_ALIGN.RIGHT)])


# ----------------------------------------------------------------------------
# SLAYTLAR
# ----------------------------------------------------------------------------
def slayt_kapak(prs):
    s = new_slide(prs, dark=True)
    rect(s, 0, 0, 0.16, H_IN, RED)                    # sol vurgu bandı
    rect(s, 0, H_IN - 0.16, W_IN, 0.16, NAVY2)
    # ajan-ağı motifi (sağ üst, ince)
    net = [(11.0, 0.9), (12.2, 1.35), (11.7, 2.25), (12.6, 2.75), (10.6, 2.0)]
    # bağlantı çizgileri (ince constellation motifi)
    import itertools
    for (a, b) in itertools.combinations(net, 2):
        conn = s.shapes.add_connector(1, Inches(a[0] + 0.11), Inches(a[1] + 0.11),
                                      Inches(b[0] + 0.11), Inches(b[1] + 0.11))
        conn.line.color.rgb = C("2A3E63")
        conn.line.width = Pt(0.75)
        _noshadow(conn)
    for (nx, ny) in net:
        oval(s, nx, ny, 0.22, 0.22, C("35507F"))
    textbox(s, 0.9, 1.55, 8.0, 0.4,
            [P([R("TEKNOFEST 2026 · YAPAY ZEKA DİL AJANLARI YARIŞMASI",
                 C("8FB0E8"), True, 12.5, SEMI)])])
    textbox(s, 0.88, 2.0, 11.0, 1.2, [P([R("AGENTRA TECH", WHITE, True, 54, SEMI)])])
    rect(s, 0.94, 3.18, 3.2, 0.05, RED)
    textbox(s, 0.9, 3.42, 11.4, 1.0,
            [P([R("Kamu Evrak ve Yazışma Süreçleri için", C("E7ECF6"), False, 21, BODY)],
               space_after=1),
             P([R("Akıllı Agent Destek Sistemi", WHITE, True, 22, SEMI)])])
    # alt bilgi çipleri
    chips = ["1. Senaryo", "11 Uzman Ajan + Orkestratör", "Offline-First", "Apache 2.0"]
    cx = 0.9
    for ch in chips:
        wch = 0.42 + len(ch) * 0.098
        pill(s, cx, 5.35, wch, 0.42, ch, NAVY2, tcolor=C("BFD0EC"), size=11)
        cx += wch + 0.22
    textbox(s, 0.9, 6.25, 11.0, 0.5,
            [P([R("Takım Tanıtım Sunumu", C("AEC2E6"), True, 13, SEMI),
                R("   ·   Temmuz 2026   ·   Takım Kaptanı: Şeyma Nur Çebi",
                  C("8296B6"), False, 12.5, BODY)])])
    return s


def slayt_bakis(prs, n):
    s = new_slide(prs)
    header(s, "Bir Bakışta", "Bir Bakışta AGENTRA TECH")
    tiles = [
        ("11", "+1", "Uzman ajan + orkestratör", PURPLE),
        ("2", "", "Zorunlu görev, tek uçtan uca akış", BLUE),
        ("500", "+", "Test — sürekli entegrasyonda yeşil", GREEN),
        ("116", "", "Etiketli sentetik kurgu evrak", ORANGE),
        ("0", "", "KVKK sızıntısı (bağımsız denetim)", RED),
    ]
    x = 0.62
    tw = 2.36
    for (big, unit, label, acc) in tiles:
        kpi_tile(s, x, 1.75, tw, 1.55, big, unit, label, acc)
        x += tw + 0.11
    # değer önerisi bandı
    rrect(s, 0.62, 3.62, 12.09, 1.15, SOFT, line=LINE, line_w=1.0, radius=0.06)
    rect(s, 0.62, 3.62, 0.09, 1.15, RED)
    textbox(s, 0.95, 3.78, 11.5, 0.9,
            [P([R("Kamu kurumlarına gelen evrağı ", INK, False, 14.5, BODY),
                R("okuyan · anlayan · eksiğini bulan · mevzuat öneren · resmî yazı taslaklayan · doğru birime yönlendiren",
                  INK, True, 14.5, SEMI),
                R(" çok ajanlı bir sistem geliştirdik.", INK, False, 14.5, BODY)], line=1.12),
             P([R("Çalışan hâli bugün elimizde — “hedefliyoruz” değil, ", MUTED, False, 12.5, BODY),
                R("çalışıyor.", RED, True, 12.5, SEMI)], space_before=3)])
    # mini iş akışı şeridi
    strip = [("📥 Girdi", C("41506B")), ("🧠 Orkestratör", PURPLE),
             ("📋 Görev 1", BLUE), ("✍️ Görev 2", ORANGE), ("📤 12+ Çıktı", GREEN)]
    x = 0.62
    sw = 2.18
    for i, (t, col) in enumerate(strip):
        node(s, x, 5.15, sw, 0.66, t, fill=col, tcolor=WHITE, line=None,
             tsize=12, radius=0.16)
        x += sw
        if i < len(strip) - 1:
            arrow(s, x - 0.06, 5.33, 0.3, 0.3, C("B7C1D4"))
            x += 0.28
    textbox(s, 0.62, 6.05, 12.0, 0.4,
            [P([R("Framework bağımsız saf Python orkestrasyon · çevrimdışı-öncelikli · isteğe bağlı LLM eskalasyonu",
                 MUTED, False, 11, BODY)], align=PP_ALIGN.CENTER)])
    footer(s, n)
    return s


def slayt_uyeler(prs, n):
    s = new_slide(prs)
    header(s, "Takım", "Takım Üyeleri ve Rolleri")
    members = [
        ("Şeyma Nur Çebi", "TAKIM KAPTANI · YAZILIM", SEYMA,
         "Görev 1 içerik analizi: sınıflandırma, bilgi çıkarımı, mevzuat RAG; değerlendirme ve uçtan uca entegrasyon.", True),
        ("Muhammed Sina Gün", "YAZILIM", SINA,
         "Mimari ve orkestrasyon; model-agnostik LLM katmanı; Görev 2 taslak üretimi (OCR ve özet dâhil).", False),
        ("Emine Elik", "VERİ · TEST · DOKÜMAN", EMINE,
         "Veri seti ve etiketleme; test kapsamı; dokümantasyon; sunum ve demo; şartname uyum takibi.", False),
        ("Zeynep Akel", "YAZILIM", ZEYNEP,
         "Görev 1 eksik bilgi tespiti; Görev 2 yönlendirme ve kullanıcı etkileşimi; triyaj, KVKK; web arayüzü.", False),
    ]
    pos = [(0.62, 1.72), (6.72, 1.72), (0.62, 3.94), (6.72, 3.94)]
    cw, ch = 6.0, 2.06
    for (name, role, col, desc, cap), (x, y) in zip(members, pos):
        rrect(s, x, y, cw, ch, WHITE, line=LINE, line_w=1.25, radius=0.07)
        rect(s, x, y, 0.11, ch, col)
        # avatar dairesi (baş harf)
        oval(s, x + 0.28, y + 0.28, 0.72, 0.72, col)
        shape_text_sp = s.shapes[-1]
        shape_text(shape_text_sp, [P([R(name[0], WHITE, True, 24, SEMI)],
                   align=PP_ALIGN.CENTER, space_after=0)], anchor=MSO_ANCHOR.MIDDLE)
        textbox(s, x + 1.18, y + 0.24, cw - 1.4, 0.4,
                [P([R(name, INK, True, 15.5, SEMI)])])
        pill(s, x + 1.18, y + 0.66, min(cw - 1.4, 0.5 + len(role) * 0.082), 0.3,
             role, col, tcolor=WHITE, size=9)
        textbox(s, x + 0.3, y + 1.14, cw - 0.55, 0.85,
                [P([R(desc, C("34405A"), False, 10.8, BODY)], line=1.05)])
        if cap:
            pill(s, x + cw - 1.15, y + 0.2, 0.95, 0.3, "★ KAPTAN", RED, size=8.5)
    textbox(s, 0.62, 6.18, 12.1, 0.6,
            [P([R("Her üye için okul/bölüm bilgisi: ", MUTED, False, 10.5, BODY),
                R("[Bölüm / Üniversite — doldurulacak]", RED, False, 10.5, BODY),
                R("      Danışman: ", MUTED, False, 10.5, BODY),
                R("[Ad Soyad — Unvan/Kurum · varsa]", RED, False, 10.5, BODY)])])
    footer(s, n)
    return s


def slayt_gorev_dagilimi(prs, n):
    s = new_slide(prs)
    header(s, "Organizasyon", "Görev Dağılımı · Şartname m.6.4")
    person_col = {"Şeyma Nur": SEYMA, "Sina": SINA, "Zeynep": ZEYNEP, "Emine": EMINE}

    def panel(x, w, title, acc, rows):
        rrect(s, x, 1.68, w, 4.15, WHITE, line=LINE, line_w=1.25, radius=0.05)
        rrect(s, x, 1.68, w, 0.5, acc, line=None, radius=0.05)
        rect(s, x, 2.02, w, 0.16, acc)  # köşe düzeltme
        shape_text(s.shapes[-2], [P([R(title, WHITE, True, 13, SEMI)],
                   align=PP_ALIGN.CENTER, space_after=0)], anchor=MSO_ANCHOR.MIDDLE)
        ry = 2.34
        for (cap, who) in rows:
            if rows.index((cap, who)) % 2 == 1:
                rect(s, x + 0.12, ry, w - 0.24, 0.5, SOFT)
            textbox(s, x + 0.28, ry + 0.11, w - 2.0, 0.32,
                    [P([R(cap, INK, False, 11.5, BODY)])])
            pw = 0.4 + len(who) * 0.095
            pill(s, x + w - pw - 0.28, ry + 0.1, pw, 0.32, who, person_col[who], size=9.5)
            ry += 0.545

    panel(0.62, 6.0, "GÖREV 1 — Sınıflandırma & İçerik Analizi", BLUE, [
        ("OCR / metin okuma", "Sina"),
        ("Tür belirleme (sınıflandırma)", "Şeyma Nur"),
        ("Bilgi çıkarımı", "Şeyma Nur"),
        ("Eksik bilgi tespiti", "Zeynep"),
        ("Mevzuat önerisi (BM25 RAG)", "Şeyma Nur"),
        ("Özet oluşturma", "Sina"),
    ])
    panel(6.72, 6.0, "GÖREV 2 — Taslaklama & Yönlendirme", ORANGE, [
        ("Resmî yazı taslağı", "Sina"),
        ("Format öz-denetimi (üslup)", "Sina"),
        ("Birim yönlendirme", "Zeynep"),
        ("Kullanıcı bilgilendirme", "Zeynep"),
        ("Eksik bilgi talebi", "Zeynep"),
    ])
    # lejant + altyapı
    textbox(s, 0.62, 6.0, 6.0, 0.4, [P([R("Her yetenek → bir sorumlu + gerçek bir kod modülü",
            MUTED, False, 10.5, BODY)])])  # placeholder, overwritten below
    lx = 6.72
    textbox(s, lx, 5.95, 0.9, 0.3, [P([R("Ekip:", MUTED, True, 10, SEMI)])])
    lx += 0.7
    for name, col in person_col.items():
        pill(s, lx, 5.93, 0.35 + len(name) * 0.093, 0.3, name, col, size=9)
        lx += 0.35 + len(name) * 0.093 + 0.15
    footer(s, n)
    return s


def slayt_mimari(prs, n):
    s = new_slide(prs)
    header(s, "Teknik Mimari", "Sistem Mimarisi — 11 Uzman Ajan + Orkestratör")
    # üst bant: Girdi -> Orkestratör -> Görev 1 lane
    node(s, 0.62, 1.9, 1.7, 1.0, "📥 Girdi", "TXT · PDF · Görüntü",
         fill=C("41506B"), tcolor=WHITE, line=None, tsize=12, ssize=8.5, radius=0.1)
    arrow(s, 2.4, 2.24, 0.34, 0.32, C("B7C1D4"))
    node(s, 2.82, 1.9, 1.66, 1.0, "🧠 Orkestratör", "koşullu akış · süre & güven",
         fill=PURPLE, tcolor=WHITE, line=None, tsize=12.5, ssize=8.5, radius=0.1)
    arrow(s, 4.56, 2.24, 0.34, 0.32, C("B7C1D4"))
    # GÖREV 1 konteyner
    rrect(s, 4.98, 1.66, 7.73, 2.02, C("EAF2FB"), line=BLUE, line_w=1.25, radius=0.05)
    textbox(s, 5.16, 1.74, 7.4, 0.3, [P([R("GÖREV 1 — Sınıflandırma & İçerik Analizi",
            BLUE, True, 11, SEMI)])])
    g1 = ["🔤 OCR", "🏷️ Sınıflandırma", "🔍 Bilgi Çıkarımı", "⚠️ Eksik Bilgi",
          "📚 Mevzuat RAG", "⏱️ Triyaj", "📝 Özet", "🔒 KVKK Maskeleme"]
    cw, ch = 1.82, 0.66
    for i, t in enumerate(g1):
        cx = 5.18 + (i % 4) * 1.87
        cy = 2.12 + (i // 4) * 0.76
        node(s, cx, cy, cw, ch, t, fill=WHITE, tcolor=C("14406E"), line=C("BBD4EC"),
             tsize=10.5, radius=0.14)
    # aşağı ok (Görev1 -> Görev2)
    arrow(s, 8.5, 3.72, 0.34, 0.42, C("B7C1D4"), direction="down")
    # alt bant: GÖREV 2 lane -> Çıktı
    rrect(s, 4.98, 4.24, 5.4, 1.62, C("FDF0E7"), line=ORANGE, line_w=1.25, radius=0.05)
    textbox(s, 5.16, 4.32, 5.1, 0.3, [P([R("GÖREV 2 — Taslaklama & Yönlendirme",
            ORANGE, True, 11, SEMI)])])
    g2 = ["📄 Taslak + Format\nÖz-Denetimi", "🏢 Birim\nYönlendirme", "💬 Kullanıcı\nBilgilendirme"]
    for i, t in enumerate(g2):
        node(s, 5.18 + i * 1.68, 4.66, 1.6, 1.02, t.replace("\n", " "),
             fill=WHITE, tcolor=C("8A3B08"), line=C("F1C9AC"), tsize=10, radius=0.1)
    arrow(s, 10.46, 4.9, 0.34, 0.34, C("B7C1D4"))
    node(s, 10.86, 4.42, 1.85, 1.26, "📤 12+ Yapılandırılmış Çıktı", None,
         fill=GREEN, tcolor=WHITE, line=None, tsize=11.5, radius=0.08)
    # açıklama şeridi
    textbox(s, 0.62, 6.12, 12.1, 0.7,
            [P([R("Her ajan tek sorumluluk üstlenir; ajanlar birbirini çağırmaz, paylaşılan durum nesnesi ",
                 MUTED, False, 11, BODY),
                R("(AgentState)", INK, True, 11, MONO),
                R(" üzerinden koordine olur. LangChain/LangGraph yok — şeffaf, hafif, denetlenebilir.",
                  MUTED, False, 11, BODY)], line=1.05)])
    footer(s, n)
    return s


def slayt_kapilar(prs, n):
    s = new_slide(prs)
    header(s, "Sorumlu Otomasyon", "Orkestratör ve 3 Koşullu Kapı")
    ymid = 2.55
    node(s, 0.62, ymid, 1.4, 0.9, "📥 Evrak", fill=C("41506B"), tcolor=WHITE,
         line=None, tsize=12, radius=0.12)
    arrow(s, 2.06, ymid + 0.28, 0.28, 0.34, C("B7C1D4"))
    diamond(s, 2.4, ymid - 0.18, 1.7, 1.26,
            [("KAPI 1", True, 10.5), ("Okunabilir mi?", False, 9), ("(≥ 30 karakter)", False, 8)])
    arrow(s, 4.16, ymid + 0.28, 0.26, 0.34, C("B7C1D4"))
    diamond(s, 4.46, ymid - 0.18, 1.7, 1.26,
            [("KAPI 2", True, 10.5), ("Metin Türkçe mi?", False, 9)])
    arrow(s, 6.22, ymid + 0.28, 0.26, 0.34, C("B7C1D4"))
    node(s, 6.5, ymid + 0.05, 1.5, 0.82, "🏷️ Sınıflandırma", fill=WHITE, tcolor=INK,
         line=LINE, tsize=11, radius=0.12)
    arrow(s, 8.04, ymid + 0.28, 0.26, 0.34, C("B7C1D4"))
    diamond(s, 8.34, ymid - 0.18, 1.7, 1.26,
            [("KAPI 3", True, 10.5), ("Güven ≥ 0,6 ?", False, 9)])
    arrow(s, 10.1, ymid + 0.28, 0.26, 0.34, GREEN)
    node(s, 10.4, ymid + 0.02, 2.3, 0.86, "✓ Akış devam eder", fill=C("E7F6EC"),
         tcolor=C("1E6B32"), line=GREEN, tsize=11.5, radius=0.12)
    # dallanmalar (aşağı)
    arrow(s, 3.1, ymid + 1.12, 0.3, 0.5, RED, direction="down")
    node(s, 2.2, ymid + 1.66, 2.1, 0.82, "⛔ Süreç durur", "uydurma çıktı üretilmez",
         fill=C("FCEBED"), tcolor=C("9E1B2C"), line=RED, tsize=11, ssize=8.5, radius=0.1)
    arrow(s, 9.05, ymid + 1.12, 0.3, 0.5, ORANGE, direction="down")
    node(s, 7.9, ymid + 1.66, 2.9, 0.82, "👤 İnsan onayı gerekli",
         "en olası 2 aday + gerekçe gösterilir",
         fill=C("FDF1E7"), tcolor=C("8A3B08"), line=ORANGE, tsize=11, ssize=8.5, radius=0.1)
    # ilke bandı
    rrect(s, 0.62, 6.02, 12.09, 0.82, NAVY, line=None, radius=0.08)
    textbox(s, 0.9, 6.14, 11.6, 0.6,
            [P([R("İlke:  ", YELLOW, True, 13, SEMI),
                R("Emin olmadığında dur, insana devret. ", WHITE, True, 13, SEMI),
                R("Kapılar halüsinasyonu önler ve insan gözetimini (HITL) mimariye kodlar — kamu gerçekliğine uygun sorumlu otomasyon.",
                  C("C7D5EC"), False, 11.5, BODY)], line=1.05)])
    footer(s, n)
    return s


def slayt_kod(prs, n):
    s = new_slide(prs)
    header(s, "Mimariye Kodlanmış İlke", "Offline-First + İnsan Onayı — Gerçek Kod")
    code_block(s, 0.62, 1.72, 6.05, 3.5, "src/agents/orchestrator.py", [
        "_MIN_ANLAMLI_KARAKTER = 30",
        "_INSAN_ONAYI_GUVEN_ESIGI = 0.6",
        "",
        "# KAPI 1 — boş/bozuk metinde analiz yapma",
        "if not self._metin_okunabilir_mi():",
        "    self._uygula_bos_metin_kapisi(mode)",
        "",
        "# KAPI 2 — Türkçe değilse taslağı durdur",
        "metin_turkce = self._metin_turkce_mi()",
        "",
        "self._run_step('classification', 'Sınıflandırma')",
        "# KAPI 3 — güven < 0.6 ise insan onayına düşür",
        "self._degerlendir_siniflandirma_guveni()",
    ], size=12)
    code_block(s, 6.83, 1.72, 5.88, 3.5, "src/agents/classification_agent.py", [
        "_ENSEMBLE_KURAL_AGIRLIGI = 0.6  # kalibre kural",
        "_ENSEMBLE_ML_AGIRLIGI    = 0.4  # saf-Python NB",
        "_ESKALASYON_ESIGI        = 0.6",
        "",
        "# Üçlü hibrit: kural + istatistik + (ops.) LLM",
        "birlesik = 0.6 * kural_ola + 0.4 * ml_ola",
        "",
        "if guven < _ESKALASYON_ESIGI and llm.aktif:",
        "    sonuc = llm_ile_dogrula(metin)",
        "# offline'da kural sonucu her zaman korunur",
        "return sonuc",
    ], size=12)
    cards = [
        ("🔌", "Offline-first", "LLM/internet olmadan tüm ajanlar kural tabanlı çalışır", GREEN),
        ("🧩", "Framework'süz", "Saf Python orkestrasyon — kara kutu yok, denetlenebilir", PURPLE),
        ("⚖️", "Mevzuat-temelli", "Her taslak kuralı yönetmelik madde/fıkra dayanağıyla denetlenir", BLUE),
    ]
    x = 0.62
    cw = 3.95
    for (ic, tt, bd, ac) in cards:
        rrect(s, x, 5.42, cw, 1.24, WHITE, line=LINE, line_w=1.25, radius=0.08)
        rect(s, x, 5.42, 0.09, 1.24, ac)
        textbox(s, x + 0.24, 5.54, cw - 0.4, 0.35,
                [P([R(ic + "  ", INK, True, 13, SEMI), R(tt, INK, True, 12.5, SEMI)])])
        textbox(s, x + 0.24, 5.92, cw - 0.42, 0.7,
                [P([R(bd, C("34405A"), False, 10.3, BODY)], line=1.05)])
        x += cw + 0.13
    footer(s, n)
    return s


def slayt_yontem(prs, n):
    s = new_slide(prs)
    header(s, "Yöntem", "Yöntemimiz ve Özgün Katkılarımız")
    # sol: hibrit sınıflandırma mini-diyagramı
    rrect(s, 0.62, 1.75, 5.5, 4.05, SOFT, line=LINE, line_w=1.0, radius=0.05)
    textbox(s, 0.85, 1.9, 5.1, 0.35, [P([R("Üçlü Hibrit Sınıflandırma", INK, True, 13, SEMI)])])
    node(s, 0.9, 2.5, 2.3, 0.8, "① Kural Skoru", "ağırlıklı kelime + yapısal",
         fill=WHITE, tcolor=BLUE, line=C("BBD4EC"), tsize=11, ssize=8, radius=0.12)
    node(s, 3.5, 2.5, 2.3, 0.8, "② Naive Bayes", "saf-Python, TF-IDF + n-gram",
         fill=WHITE, tcolor=TEAL, line=C("A9E3D2"), tsize=11, ssize=8, radius=0.12)
    arrow(s, 1.7, 3.35, 0.3, 0.34, C("B7C1D4"), direction="down")
    arrow(s, 4.3, 3.35, 0.3, 0.34, C("B7C1D4"), direction="down")
    node(s, 1.7, 3.78, 2.9, 0.72, "Ensemble  0,6 × kural + 0,4 × ML",
         fill=PURPLE, tcolor=WHITE, line=None, tsize=11, radius=0.12)
    arrow(s, 2.95, 4.52, 0.3, 0.3, C("B7C1D4"), direction="down")
    node(s, 1.5, 4.86, 3.3, 0.72, "③ LLM Eskalasyonu (güven < 0,6)",
         "yoksa kural sonucu korunur",
         fill=WHITE, tcolor=ORANGE, line=C("F1C9AC"), tsize=10.5, ssize=8, radius=0.1)
    # sağ: özgünlük kartları
    items = [
        ("🧩", "Framework bağımsız saf Python çok ajanlı orkestrasyon",
         "Kara kutu agent kütüphanesi yok; her karar izlenebilir.", PURPLE),
        ("✅", "Taslak format öz-denetimi",
         "Sistem kendi çıktısını Resmî Yazışma Yönetmeliği'nin kurallarına göre puanlar.", BLUE),
        ("🔁", "Offline-first hibrit tasarım",
         "LLM'siz tam işlev + LLM'li eskalasyon aynı mimaride.", GREEN),
        ("🛡️", "Düşük güven kapısı + insan onayı",
         "Halüsinasyon önleme ve HITL kamu gerçekliğine uygun.", RED),
    ]
    y = 1.75
    for (ic, tt, bd, ac) in items:
        rrect(s, 6.32, y, 6.39, 0.94, WHITE, line=LINE, line_w=1.25, radius=0.08)
        rect(s, 6.32, y, 0.09, 0.94, ac)
        textbox(s, 6.58, y + 0.12, 6.0, 0.35,
                [P([R(ic + "  ", INK, True, 13, SEMI), R(tt, INK, True, 12.5, SEMI)])])
        textbox(s, 6.58, y + 0.5, 6.0, 0.4,
                [P([R(bd, C("34405A"), False, 10.5, BODY)], line=1.0)])
        y += 1.02
    footer(s, n)
    return s


def slayt_metrikler(prs, n):
    s = new_slide(prs)
    header(s, "Uygulama", "Çalışan Sistem — Ölçülebilir Başarım")

    def panel(x, w, title, acc, bars, extra):
        rrect(s, x, 1.72, w, 3.65, WHITE, line=LINE, line_w=1.25, radius=0.05)
        rrect(s, x, 1.72, w, 0.5, acc, line=None, radius=0.05)
        rect(s, x, 2.06, w, 0.16, acc)
        shape_text(s.shapes[-2], [P([R(title, WHITE, True, 12.5, SEMI)],
                   align=PP_ALIGN.CENTER, space_after=0)], anchor=MSO_ANCHOR.MIDDLE)
        by = 2.42
        for (lbl, val) in bars:
            metric_bar(s, x + 0.35, by, w - 0.7, lbl, val, acc)
            by += 0.62
        textbox(s, x + 0.35, by + 0.02, w - 0.7, 0.4,
                [P([R(extra, C("34405A"), False, 10.5, BODY)])])

    panel(0.62, 5.95, "Geliştirme Seti · 52 evrak", BLUE, [
        ("Sınıflandırma doğruluğu", 1.00),
        ("Birim yönlendirme", 0.96),
        ("Eksik bilgi tespiti (F1)", 1.00),
        ("Mevzuat isabet@3", 0.96),
    ], "Taslak kalitesi 93,6/100   ·   KVKK sızıntısı 0")
    panel(6.77, 5.95, "Adversarial Tutulmuş Set · 16 evrak — hiç dokunulmamış", ORANGE, [
        ("Sınıflandırma doğruluğu", 0.94),
        ("Birim yönlendirme", 1.00),
        ("Eksik bilgi tespiti (F1)", 0.83),
        ("Mevzuat isabet@3", 0.94),
    ], "Taslak kalitesi 95,8/100   ·   KVKK sızıntısı 0")
    rrect(s, 0.62, 5.55, 12.09, 1.02, SOFT, line=LINE, line_w=1.0, radius=0.06)
    textbox(s, 0.9, 5.66, 11.7, 0.85,
            [P([R("Kaynak: ", MUTED, True, 10.5, SEMI),
                R("data/processed/eval_report*.json", INK, False, 10.5, MONO),
                R("  ·  tamamen çevrimdışı mod  ·  500+ test yeşil  ·  evrak başına saniye-altı işleme (medyan ~0,1 sn)",
                  MUTED, False, 10.5, BODY)], line=1.05),
             P([R("Dürüstlük: ", RED, True, 10.5, SEMI),
                R("Adversarial sette eksik-bilgi F1 (0,83) düşüşü tuzakların hedefidir; sonuçlar hiçbir düzeltme yapılmadan olduğu gibi raporlanmıştır.",
                  MUTED, False, 10.5, BODY)], space_before=2, line=1.02)])
    footer(s, n)
    return s


def slayt_motivasyon(prs, n):
    s = new_slide(prs)
    header(s, "Neden Bu Proje?", "Katılım Motivasyonumuz ve Gerekçeleri")
    items = [
        ("⏱️", "Ölçülebilir verimlilik",
         "İlk inceleme, taslak ve yönlendirme adımlarında personel zamanı kazanımı.", BLUE),
        ("🔒", "Veri egemenliği & KVKK",
         "Kişisel veriyi 3. taraf API'ye sızdırmadan, tamamen yerel (offline) çalışan yerli çözüm.", RED),
        ("🌍", "Açık kaynak katkısı",
         "Türkçe dil teknolojileri ekosistemine (TAKP) Apache 2.0 lisanslı, tekrar-üretilebilir katkı.", GREEN),
        ("🧭", "Sorumlu otomasyon",
         "Emin olmadığı kararda durup insana devreden, gerekçeli ve denetlenebilir tasarım.", PURPLE),
        ("🎓", "Takımın ilgisi ve birikimi",
         "Türkçe NLP ve kamu süreçlerine dair merak ve ortak çalışma isteği.", ORANGE),
    ]
    pos = [(0.62, 1.78), (4.66, 1.78), (8.7, 1.78), (0.62, 4.0), (4.66, 4.0)]
    cw, ch = 3.9, 2.02
    for (ic, tt, bd, ac), (x, y) in zip(items, pos):
        rrect(s, x, y, cw, ch, WHITE, line=LINE, line_w=1.25, radius=0.08)
        oval(s, x + 0.26, y + 0.26, 0.66, 0.66, C("EEF3FB"))
        shape_text(s.shapes[-1], [P([R(ic, INK, False, 20, BODY)],
                   align=PP_ALIGN.CENTER, space_after=0)], anchor=MSO_ANCHOR.MIDDLE)
        rect(s, x + 0.26, y + ch - 0.34, 0.66, 0.06, ac)
        textbox(s, x + 1.06, y + 0.34, cw - 1.25, 0.55,
                [P([R(tt, INK, True, 13, SEMI)], line=1.0)])
        textbox(s, x + 0.28, y + 1.08, cw - 0.5, 0.85,
                [P([R(bd, C("34405A"), False, 10.6, BODY)], line=1.06)])
    # büyük vurgu kutusu (5. hücre yanı)
    rrect(s, 8.7, 4.0, 3.9, 2.02, NAVY, line=None, radius=0.08)
    textbox(s, 8.98, 4.2, 3.4, 1.7,
            [P([R("Kamuya uygun yapay zekâ", YELLOW, True, 14, SEMI)], space_after=4),
             P([R("“Her şeyi otomatikleştiren” değil; ", C("D8E2F3"), False, 11.5, BODY),
                R("doğru yerde insana devreden", WHITE, True, 11.5, SEMI),
                R(" bir sistem güven verir.", C("D8E2F3"), False, 11.5, BODY)], line=1.1)])
    footer(s, n)
    return s


def slayt_acik_kaynak(prs, n):
    s = new_slide(prs)
    header(s, "İlkeler", "Açık Kaynak · KVKK · Dürüstlük")
    card(s, 0.62, 1.8, 3.9, 3.5, "📖", "Açık Kaynak",
         ["Apache 2.0 lisansı", "Tüm dokümantasyon Türkçe",
          "Model ağırlığı depoya yüklenmez (yalnızca bağlantı+sürüm+lisans)",
          "TAKP GitHub uyumu"], GREEN)
    card(s, 4.66, 1.8, 3.9, 3.5, "🔒", "KVKK",
         ["9 kategori kişisel veri format-koruyarak maskelenir",
          "Sızıntı bağımsız denetçiyle ölçülür",
          "Sızıntısız oran 1,00 (her sette 0 kaçak)",
          "Tamamen kural tabanlı, offline"], BLUE)
    card(s, 8.7, 1.8, 3.99, 3.5, "🎯", "Dürüstlük",
         ["Ölçümler ne çıkarsa olduğu gibi raporlanır",
          "Tutulmuş setlerdeki hatalar gizlenmez",
          "Gerçek kamu verisi ASLA kullanılmaz",
          "Yalnızca sentetik/kurgu evrak + kamuya açık mevzuat"], RED)
    rrect(s, 0.62, 5.5, 12.09, 1.05, SOFT, line=LINE, line_w=1.0, radius=0.06)
    rect(s, 0.62, 5.5, 0.09, 1.05, INK)
    textbox(s, 0.95, 5.64, 11.6, 0.85,
            [P([R("Depo:  ", MUTED, True, 11.5, SEMI),
                R("github.com/msgxr/teknofest-2026-kamu-evrak-akilli-ajan", INK, True, 11.5, MONO)]),
             P([R("116 etiketli sentetik evrak (52 geliştirme + 64 tutulmuş)  ·  15 belgelik mevzuat korpusu  ·  508+ test",
                 MUTED, False, 11, BODY)], space_before=3)])
    footer(s, n)
    return s


def slayt_yol_haritasi(prs, n):
    s = new_slide(prs)
    header(s, "Yol Haritası", "Bugünden Finale ve Ürünleşmeye")
    steps = [
        ("BUGÜN", "Çalışan sistem", "CLI + Streamlit + demo; iki görev tek uçtan uca akışta; 508+ test yeşil", GREEN),
        ("FİNALE KADAR", "Derinleştirme", "Türkçe NER · semantik katman kalibrasyonu · kullanıcı testleri ve dürüst sürekli ölçüm", BLUE),
        ("ÜRÜNLEŞME", "Sahaya çıkış", "EBYS (TS 13298) / e-Yazışma entegrasyonu · pilot kurum · yerel model = veri kurumdan çıkmaz", ORANGE),
    ]
    x = 0.62
    cw = 3.9
    for i, (kick, tt, bd, ac) in enumerate(steps):
        rrect(s, x, 2.1, cw, 2.5, WHITE, line=LINE, line_w=1.25, radius=0.07)
        rrect(s, x, 2.1, cw, 0.56, ac, line=None, radius=0.07)
        rect(s, x, 2.42, cw, 0.24, ac)
        textbox(s, x, 2.12, cw, 0.5,
                [P([R(kick, WHITE, True, 12, SEMI)], align=PP_ALIGN.CENTER)],
                anchor=MSO_ANCHOR.MIDDLE)
        textbox(s, x + 0.28, 2.86, cw - 0.5, 0.5, [P([R(tt, INK, True, 15, SEMI)])])
        textbox(s, x + 0.28, 3.4, cw - 0.52, 1.1,
                [P([R(bd, C("34405A"), False, 11, BODY)], line=1.12)])
        x += cw
        if i < len(steps) - 1:
            arrow(s, x - 0.04, 3.15, 0.32, 0.36, C("AEB9CE"))
            x += 0.29
    # yarışma takvimi şeridi
    rrect(s, 0.62, 4.95, 12.09, 1.55, NAVY, line=None, radius=0.06)
    textbox(s, 0.9, 5.08, 11.5, 0.3,
            [P([R("YARIŞMA TAKVİMİ ", YELLOW, True, 11, SEMI),
                R("(teknofest.org'dan teyit edilmeli)", C("8FA6C6"), False, 9.5, BODY)])])
    dates = [("17 Tem", "Son teslim"), ("19 Tem", "Ön değ. sonucu"),
             ("21 Tem", "Teknik sınav"), ("24 Tem", "Finalistler"), ("Ağustos", "Final")]
    dx = 0.9
    dw = 2.3
    for i, (d, lbl) in enumerate(dates):
        oval(s, dx, 5.5, 0.16, 0.16, YELLOW)
        textbox(s, dx + 0.24, 5.42, dw, 0.3, [P([R(d, WHITE, True, 12, SEMI)])])
        textbox(s, dx + 0.24, 5.74, dw, 0.3, [P([R(lbl, C("B9C8E2"), False, 10, BODY)])])
        if i < len(dates) - 1:
            rect(s, dx + 0.16, 5.575, dw - 0.02, 0.02, C("2E4670"))
        dx += dw
    footer(s, n)
    return s


def slayt_kapanis(prs, n):
    s = new_slide(prs, dark=True)
    rect(s, 0, 0, 0.16, H_IN, RED)
    textbox(s, 0.9, 1.15, 11.0, 0.4,
            [P([R("TEŞEKKÜRLER", RED, True, 15, SEMI)])])
    textbox(s, 0.88, 1.55, 11.2, 1.0, [P([R("AGENTRA TECH", WHITE, True, 42, SEMI)])])
    textbox(s, 0.9, 2.55, 11.4, 0.5,
            [P([R("Kamu Evrak ve Yazışma Süreçleri için Akıllı Agent Destek Sistemi",
                 C("D8E2F3"), False, 15, BODY)])])
    # üye şeridi
    members = [("Şeyma Nur Çebi", "Kaptan", SEYMA), ("Muhammed Sina Gün", "", SINA),
               ("Emine Elik", "", EMINE), ("Zeynep Akel", "", ZEYNEP)]
    x = 0.9
    for (nm, rl, col) in members:
        w = 0.5 + len(nm) * 0.11 + (len(rl) * 0.09 if rl else 0)
        rrect(s, x, 3.35, w, 0.55, NAVY2, line=None, radius=0.3)
        oval(s, x + 0.16, 3.5, 0.25, 0.25, col)
        runs = [R(nm, WHITE, True, 11.5, SEMI)]
        if rl:
            runs.append(R("  · " + rl, C("9FB4D8"), False, 10, BODY))
        textbox(s, x + 0.5, 3.46, w - 0.5, 0.35, [P(runs)])
        x += w + 0.2
    # iletişim kutuları
    rrect(s, 0.9, 4.35, 7.4, 0.7, NAVY2, line=None, radius=0.1)
    textbox(s, 1.15, 4.5, 7.0, 0.4,
            [P([R("Depo:  ", C("9FB4D8"), True, 12, SEMI),
                R("github.com/msgxr/teknofest-2026-kamu-evrak-akilli-ajan", WHITE, False, 12, MONO)])])
    rrect(s, 0.9, 5.2, 7.4, 0.7, NAVY2, line=None, radius=0.1)
    textbox(s, 1.15, 5.35, 7.0, 0.4,
            [P([R("İletişim:  ", C("9FB4D8"), True, 12, SEMI),
                R("[takım e-postası — doldurulacak]", C("F3B7C0"), False, 12, BODY)])])
    textbox(s, 0.9, 6.35, 11.4, 0.6,
            [P([R("Kamuya uygun, dürüst, açık kaynak ve ", C("C7D5EC"), False, 13.5, BODY),
                R("bugün çalışan", WHITE, True, 13.5, SEMI),
                R(" bir sistem sunuyoruz.", C("C7D5EC"), False, 13.5, BODY)])])
    return s


def uret(cikti: Path) -> int:
    prs = Presentation()
    prs.slide_width = Inches(W_IN)
    prs.slide_height = Inches(H_IN)
    prs.core_properties.author = "AGENTRA TECH"
    prs.core_properties.last_modified_by = "AGENTRA TECH"
    prs.core_properties.title = "AGENTRA TECH — Takım Tanıtım Sunumu (TEKNOFEST 2026 TYDA · 1. Senaryo)"

    slayt_kapak(prs)
    slayt_bakis(prs, 2)
    slayt_uyeler(prs, 3)
    slayt_gorev_dagilimi(prs, 4)
    slayt_mimari(prs, 5)
    slayt_kapilar(prs, 6)
    slayt_kod(prs, 7)
    slayt_yontem(prs, 8)
    slayt_metrikler(prs, 9)
    slayt_motivasyon(prs, 10)
    slayt_acik_kaynak(prs, 11)
    slayt_yol_haritasi(prs, 12)
    slayt_kapanis(prs, 13)

    cikti.parent.mkdir(parents=True, exist_ok=True)
    try:
        prs.save(str(cikti))
        return len(prs.slides._sldIdLst), cikti
    except PermissionError:
        # Hedef dosya açık/kilitli (PowerPoint'te açık olabilir) — yedek ada yaz
        yedek = cikti.with_name(cikti.stem + "_yeni.pptx")
        prs.save(str(yedek))
        return len(prs.slides._sldIdLst), yedek


def main() -> None:
    ap = argparse.ArgumentParser(description="AGENTRA TECH Takım Tanıtım Sunumu üreticisi")
    ap.add_argument("--cikti", type=Path,
                    default=PROJE_KOKU / "presentations" / "Agentra_Tech_Takim_Tanitim_Sunum.pptx")
    args = ap.parse_args()
    adet, yazilan = uret(args.cikti)
    if yazilan != args.cikti:
        print(f"UYARI: {args.cikti.name} acik/kilitli oldugu icin yedek ada yazildi.")
    print(f"Sunum uretildi: {yazilan} ({adet} slayt)")
    print("Hatirlatma: PDF surumunu PowerPoint'ten 'Farkli Kaydet -> PDF' ile alin.")


if __name__ == "__main__":
    main()
